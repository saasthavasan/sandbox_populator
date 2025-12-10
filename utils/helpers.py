#!/usr/bin/env python3
"""
Helper utilities for file creation and manipulation
"""

import os
import random
import string
from pathlib import Path
from datetime import datetime, timedelta


def ensure_directory(path):
    """Create directory if it doesn't exist"""
    Path(path).mkdir(parents=True, exist_ok=True)
    return path


def random_string(length=10, include_special=False):
    """Generate a random string"""
    chars = string.ascii_letters + string.digits
    if include_special:
        chars += "!@#$%^&*()"
    return ''.join(random.choice(chars) for _ in range(length))


def random_date(start_date, end_date):
    """Generate a random date between start_date and end_date"""
    time_between = end_date - start_date
    days_between = time_between.days
    random_days = random.randrange(days_between)
    return start_date + timedelta(days=random_days)


def format_currency(amount):
    """Format number as currency"""
    return f"${amount:,.2f}"


def generate_fake_ssn():
    """Generate a fake SSN (for sandbox use only)"""
    area = random.randint(100, 899)
    group = random.randint(10, 99)
    serial = random.randint(1000, 9999)
    return f"{area}-{group}-{serial}"


def generate_fake_phone():
    """Generate a fake phone number"""
    area = random.randint(200, 999)
    prefix = random.randint(200, 999)
    line = random.randint(1000, 9999)
    return f"({area}) {prefix}-{line}"


def generate_fake_email(name, domain="example.com"):
    """Generate a fake email address"""
    name = name.lower().replace(" ", ".")
    return f"{name}@{domain}"


def random_ip_address():
    """Generate a random IP address"""
    return f"{random.randint(1, 255)}.{random.randint(0, 255)}.{random.randint(0, 255)}.{random.randint(1, 255)}"


def random_mac_address():
    """Generate a random MAC address"""
    return ":".join([f"{random.randint(0, 255):02x}" for _ in range(6)])


def file_size_string(size_bytes):
    """Convert bytes to human-readable string"""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.1f} TB"


def generate_guid():
    """Generate a GUID-like string"""
    return f"{random_string(8)}-{random_string(4)}-{random_string(4)}-{random_string(4)}-{random_string(12)}"


def weighted_random_choice(choices, weights):
    """Select random item with weights"""
    return random.choices(choices, weights=weights, k=1)[0]


def create_file_safely(filepath, content, encoding='utf-8'):
    """Create a file with proper error handling"""
    try:
        ensure_directory(os.path.dirname(filepath))
        Path(filepath).write_text(content, encoding=encoding)
        return True
    except Exception as e:
        print(f"Error creating file {filepath}: {e}")
        return False


def write_text_file(path, content, encoding='utf-8'):
    """
    Write text content to a file and ensure its parent directory exists.
    Returns the created Path for convenience.
    """
    file_path = Path(path)
    ensure_directory(file_path.parent)
    file_path.write_text(content, encoding=encoding)
    return file_path


def write_binary_file(path, data: bytes):
    """
    Write binary data to a file and ensure its parent directory exists.
    Returns the created Path for convenience.
    """
    file_path = Path(path)
    ensure_directory(file_path.parent)
    file_path.write_bytes(data)
    return file_path


def chrome_timestamp(dt: datetime) -> int:
    """
    Convert a datetime to Chrome/Webkit timestamp (microseconds since Jan 1, 1601 UTC).
    """
    epoch_start = datetime(1601, 1, 1)
    delta = dt - epoch_start
    return int(delta.total_seconds() * 1_000_000)


def firefox_timestamp(dt: datetime) -> int:
    """
    Convert a datetime to Firefox timestamp (microseconds since Unix epoch).
    """
    return int(dt.timestamp() * 1_000_000)


def create_pdf(path, title, sections):
    """
    Create a simple multi-section PDF using fpdf2.
    sections: iterable of (heading, body_text)
    """
    from fpdf import FPDF
    import textwrap

    def _safe(text: str) -> str:
        # Replace characters that can't be represented to avoid PDF errors
        return text.encode("latin-1", "replace").decode("latin-1")

    def _emit_wrapped_lines(pdf_obj, text, width=90):
        """
        Write text to the PDF with manual wrapping to avoid long unbroken lines.
        """
        text = _safe(text)
        # If the line has no spaces, chunk it manually
        if " " not in text and len(text) > width:
            chunks = [text[i:i + width] for i in range(0, len(text), width)]
        else:
            chunks = textwrap.wrap(text, width=width) or [text]

        # Use explicit width to avoid zero-width errors from FPDF
        avail_width = max(int(pdf_obj.w - pdf_obj.l_margin - pdf_obj.r_margin), 20)
        for chunk in chunks:
            try:
                pdf_obj.multi_cell(avail_width, 6, chunk)
            except Exception:
                # As a last resort, shrink further
                pdf_obj.multi_cell(max(avail_width - 20, 20), 6, chunk)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", "B", 16)
    pdf.cell(0, 10, _safe(title), ln=1)
    pdf.ln(4)

    for heading, body in sections:
        pdf.set_font("Arial", "B", 12)
        pdf.cell(0, 8, _safe(heading), ln=1)
        pdf.set_font("Arial", "", 10)
        for line in body.splitlines():
            _emit_wrapped_lines(pdf, line, width=90)
        pdf.ln(2)

    pdf.output(str(path))
    return Path(path)


def create_workbook(path, sheets):
    """
    Create an XLSX workbook using openpyxl.
    sheets: list of tuples (sheet_name, rows) where rows is a list of row lists.
    """
    from openpyxl import Workbook

    wb = Workbook()
    # clear the default sheet
    wb.remove(wb.active)
    for sheet_name, rows in sheets:
        ws = wb.create_sheet(sheet_name[:31] or "Sheet")
        for row in rows:
            ws.append(row)
    wb.save(path)
    return Path(path)


def create_presentation(path, slides):
    """
    Create a PPTX presentation using python-pptx.
    slides: list of dicts {title: str, bullets: [str]}
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    for slide in slides:
        layout = prs.slide_layouts[1]  # Title and content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get("title", "")
        body = s.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(slide.get("bullets", [])):
            if i == 0:
                tf.text = bullet
                tf.paragraphs[0].font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
    prs.save(path)
    return Path(path)


def get_windows_paths(base_override: str | Path | None = None):
    """
    Resolve Windows-style user directories (Desktop, Documents, Downloads, AppData, Program Files).
    Falls back to Path.home() when environment variables are missing.
    """
    base_home = (
        Path(base_override)
        if base_override
        else Path(os.environ.get("USERPROFILE", Path.home()))
    )

    appdata_local = Path(os.environ.get("LOCALAPPDATA", base_home / "AppData" / "Local"))
    appdata_roaming = Path(os.environ.get("APPDATA", base_home / "AppData" / "Roaming"))
    program_files = Path(os.environ.get("ProgramFiles", base_home / "Program Files"))

    return {
        "home": base_home,
        "desktop": base_home / "Desktop",
        "documents": base_home / "Documents",
        "downloads": base_home / "Downloads",
        "pictures": base_home / "Pictures",
        "music": base_home / "Music",
        "videos": base_home / "Videos",
        "appdata_local": appdata_local,
        "appdata_roaming": appdata_roaming,
        "program_files": program_files,
    }
